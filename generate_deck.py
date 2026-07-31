import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Set 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    bg_color = RGBColor(10, 15, 30)         # Dark Navy
    primary_color = RGBColor(139, 92, 246)  # Electric Violet
    accent_color = RGBColor(16, 185, 129)   # Emerald Green
    text_light = RGBColor(241, 245, 249)    # Light gray/white
    text_muted = RGBColor(148, 163, 184)    # Muted gray
    box_bg = RGBColor(20, 27, 45)           # Lighter box container

    # Helper function to apply dark background
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Helper function to create title on slides
    def add_slide_header(slide, title_text, category_text=""):
        # Category Tracker (Top-left)
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
            cat_tf = cat_box.text_frame
            cat_tf.word_wrap = True
            cat_p = cat_tf.paragraphs[0]
            cat_p.text = category_text.upper()
            cat_p.font.name = 'Arial'
            cat_p.font.size = Pt(10)
            cat_p.font.bold = True
            cat_p.font.color.rgb = accent_color
            
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = 'Arial'
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = text_light

    blank_slide_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Subtitle Category tag
    tag_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.4))
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = "SOCF 2.0 GRAND FINALE  •  AGENTIC AI TRACK"
    tag_p.font.name = 'Arial'
    tag_p.font.size = Pt(11)
    tag_p.font.bold = True
    tag_p.font.color.rgb = accent_color

    # Project Title
    main_title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.8))
    mt_tf = main_title_box.text_frame
    mt_p = mt_tf.paragraphs[0]
    mt_p.text = "PaperPilot ✈️"
    mt_p.font.name = 'Arial'
    mt_p.font.size = Pt(56)
    mt_p.font.bold = True
    mt_p.font.color.rgb = text_light

    # Add color-highlighted subtitle text in the same frame
    mt_p2 = mt_tf.add_paragraph()
    mt_p2.text = "Autonomous Research Briefing & Pedagogy Agent"
    mt_p2.font.name = 'Arial'
    mt_p2.font.size = Pt(20)
    mt_p2.font.color.rgb = primary_color
    mt_p2.space_before = Pt(10)

    # Footer/Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.0))
    meta_tf = meta_box.text_frame
    meta_p = meta_tf.paragraphs[0]
    meta_p.text = "Team: Neuro Nexus  •  Team Lead: Shivanesh V  •  Agentic AI Track"
    meta_p.font.name = 'Arial'
    meta_p.font.size = Pt(12)
    meta_p.font.color.rgb = text_muted
    meta_p.font.bold = True

    # ==========================================
    # SLIDE 2: The Core Problem in Academic RAG
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "The Current Academic AI Bottlenecks", "Problem Statement")

    # Column 1: Multi-Agent Bloat
    col1 = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5))
    c1_tf = col1.text_frame
    c1_tf.word_wrap = True
    p1 = c1_tf.paragraphs[0]
    p1.text = "1. Multi-Agent Bloat"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = primary_color
    p1_sub = c1_tf.add_paragraph()
    p1_sub.text = "Traditional research frameworks perform 10+ separate LLM calls to build briefs, flashcards, and maps. This causes massive API token burn and excessive costs."
    p1_sub.font.size = Pt(12)
    p1_sub.font.color.rgb = text_muted
    p1_sub.space_before = Pt(8)

    # Column 2: Hallucination Risk
    col2 = slide2.shapes.add_textbox(Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5))
    c2_tf = col2.text_frame
    c2_tf.word_wrap = True
    p2 = c2_tf.paragraphs[0]
    p2.text = "2. Untrusted Hallucinations"
    p2.font.bold = True
    p2.font.size = Pt(18)
    p2.font.color.rgb = primary_color
    p2_sub = c2_tf.add_paragraph()
    p2_sub.text = "Most RAG search tools lack direct connection between generated briefs and original PDF locations. Users cannot quickly cross-check or verify output facts."
    p2_sub.font.size = Pt(12)
    p2_sub.font.color.rgb = text_muted
    p2_sub.space_before = Pt(8)

    # Column 3: Language Barrier
    col3 = slide2.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.5))
    c3_tf = col3.text_frame
    c3_tf.word_wrap = True
    p3 = c3_tf.paragraphs[0]
    p3.text = "3. Accessibility Barriers"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = primary_color
    p3_sub = c3_tf.add_paragraph()
    p3_sub.text = "Complex academic concepts are gatekept by dense terminology and English-only outputs. Translating via standard cloud translation APIs incurs massive subscription costs."
    p3_sub.font.size = Pt(12)
    p3_sub.font.color.rgb = text_muted
    p3_sub.space_before = Pt(8)

    # ==========================================
    # SLIDE 3: Tier 1 - Rubric Max-Out (Efficiency)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "Maxing the Rubric: Architecture & Efficiency", "Tier 1 Features")

    # Grid of 4 boxes representing Tier 1 Features
    def add_feature_box(slide, x, y, w, h, title, desc, efficiency_point):
        # Create background shape
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = box_bg
        box.line.color.rgb = RGBColor(30, 41, 59)
        
        # Text Frame
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = primary_color
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = text_light
        p_desc.space_before = Pt(5)
        
        p_eff = tf.add_paragraph()
        p_eff.text = f"⚡ {efficiency_point}"
        p_eff.font.bold = True
        p_eff.font.size = Pt(10)
        p_eff.font.color.rgb = accent_color
        p_eff.space_before = Pt(8)

    # 1. Zero-API Ingestion
    add_feature_box(slide3, 0.8, 1.8, 5.6, 2.3, 
                    "Local Ingestion Pipeline", 
                    "Chunks text page-by-page and computes sentence embeddings using local model (all-MiniLM-L6-v2) stored in ChromaDB.",
                    "0 API calls during ingestion")

    # 2. Single-Pass Master Agent
    add_feature_box(slide3, 6.8, 1.8, 5.6, 2.3, 
                    "Single-Pass JSON Synthesis", 
                    "Fetches Brief, Flashcards, Concept Map, Skeptic Critiques, and Podcast script in exactly one API call.",
                    "Reduces LLM API usage by >70%")

    # 3. SQLite Semantic Cache
    add_feature_box(slide3, 0.8, 4.5, 5.6, 2.3, 
                    "SQLite Semantic Cache", 
                    "Hashes PDF files and caching queries; repeats load instantly from database with zero backend operations.",
                    "0 API calls on Cache Hits")

    # 4. RAG Efficiency Hub
    add_feature_box(slide3, 6.8, 4.5, 5.6, 2.3, 
                    "Live 'RAG Efficiency Hub'", 
                    "A glassmorphism widgets displaying live counters for API Calls, Tokens, Cache hits, and Saved Dollars.",
                    "Proves efficiency on-screen to judges")

    # ==========================================
    # SLIDE 4: Tier 2 - The Showstopper Innovations
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "Judge Eye-Candy: Immersive Features", "Tier 2 Features")

    # 1. Local Multilingual Voice Reader
    add_feature_box(slide4, 0.8, 1.8, 5.6, 2.3, 
                    "Zero-API Multilingual Voice Reader", 
                    "Narrates summaries and scripts in Hindi, Tamil, Telugu, Marathi, Bengali, and English utilizing native Web Speech API.",
                    "0 API cost for localized speech")

    # 2. Agentic Skeptic Engine
    add_feature_box(slide4, 6.8, 1.8, 5.6, 2.3, 
                    "The 'Agentic Skeptic' Engine", 
                    "Evaluates methodology rigor (1-10) and reports the paper's biggest weakness or bias in one clear sentence.",
                    "Real agentic analysis (not just summary)")

    # 3. Paper-to-Podcast Mode
    add_feature_box(slide4, 0.8, 4.5, 5.6, 2.3, 
                    "Audio 'Paper-to-Podcast' Mode", 
                    "Generates a 2-person conversational script between a Host and Researcher explaining the paper.",
                    "Democratizes auditory learning")

    # 4. ELI5 vs. Academic Toggle
    add_feature_box(slide4, 6.8, 4.5, 5.6, 2.3, 
                    "Dynamic 'ELI5' Switcher", 
                    "Instantly switches complex academic brief summaries to simple real-world child analogies.",
                    "Extends usage to non-technical users")

    # ==========================================
    # SLIDE 5: Tier 3 - Grounding & Actionability
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "Grounding, Citations & Developer Actions", "Tier 3 Features")

    # Feature 9 Box
    box9 = slide5.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    box9.fill.solid()
    box9.fill.fore_color.rgb = box_bg
    box9.line.color.rgb = RGBColor(30, 41, 59)
    tf9 = box9.text_frame
    tf9.word_wrap = True
    p9 = tf9.paragraphs[0]
    p9.text = "9. Interactive Citation-Highlighting PDF"
    p9.font.bold = True
    p9.font.size = Pt(18)
    p9.font.color.rgb = primary_color
    p9_desc = tf9.add_paragraph()
    p9_desc.text = "Clicking any [1] citation link in the research brief dynamically updates the PDF iframe on the screen, scrolling and highlighting the exact line of the source PDF. This guarantees a zero-hallucination workflow and builds total trust."
    p9_desc.font.size = Pt(12)
    p9_desc.font.color.rgb = text_light
    p9_desc.space_before = Pt(15)

    # Feature 10 Box
    box10 = slide5.shapes.add_shape(1, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5))
    box10.fill.solid()
    box10.fill.fore_color.rgb = box_bg
    box10.line.color.rgb = RGBColor(30, 41, 59)
    tf10 = box10.text_frame
    tf10.word_wrap = True
    p10 = tf10.paragraphs[0]
    p10.text = "10. One-Click Replication Deep-Linker"
    p10.font.bold = True
    p10.font.size = Pt(18)
    p10.font.color.rgb = primary_color
    p10_desc = tf10.add_paragraph()
    p10_desc.text = "The agent extracts open-source libraries, neural models, algorithms, and datasets mentioned in the text. It automatically compiles direct GitHub repository code and Kaggle dataset query links, turning academic reading into immediate developer execution."
    p10_desc.font.size = Pt(12)
    p10_desc.font.color.rgb = text_light
    p10_desc.space_before = Pt(15)

    # ==========================================
    # SLIDE 6: 3-Minute Grand Finale Pitch Script
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "The Winning 3-Minute Finale Pitch Script", "Hackathon Pitch Formula")

    # Table layout for Pitch Timeline
    rows, cols = 5, 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(1.5)  # Time
    table.columns[1].width = Inches(2.5)  # Phase
    table.columns[2].width = Inches(7.7)  # Action & Narrative

    headers = ["TIME", "PHASE", "PITCH ACTION & NARRATIVE"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = box_bg
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = accent_color

    pitch_data = [
        ("0:00 - 0:40", "Problem & Hook", "Hook: 'Most AI research assistants burn through dozens of API calls, hallucinate, and exclude non-English speakers. We built PaperPilot specifically to solve this.'"),
        ("0:40 - 1:20", "The Efficiency Proof", "Upload PDF: 'Ingestion is 100% local (0 API calls). In a single-pass JSON call, our Master Agent fetches the Brief, 3D Flashcards, Concept Map, and Podcast dialogue.'"),
        ("1:20 - 2:20", "The Innovation Flex", "Click Skeptic Score & TTS: 'PaperPilot critiques methodology rigor. And to democratize science, local browser SpeechSynthesis reads papers in Hindi or Tamil with 0 API costs.'"),
        ("2:20 - 3:00", "The Mic Drop", "Re-upload: 'Our SQLite Semantic Cache returns cached papers instantly. Maximum API efficiency, zero hallucinations, total localized accessibility.'")
    ]

    for r_idx, (t_val, p_val, desc_val) in enumerate(pitch_data, start=1):
        # Time
        c0 = table.cell(r_idx, 0)
        c0.text_frame.paragraphs[0].text = t_val
        c0.text_frame.paragraphs[0].font.bold = True
        c0.text_frame.paragraphs[0].font.size = Pt(11)
        c0.text_frame.paragraphs[0].font.color.rgb = text_light
        
        # Phase
        c1 = table.cell(r_idx, 1)
        c1.text_frame.paragraphs[0].text = p_val
        c1.text_frame.paragraphs[0].font.bold = True
        c1.text_frame.paragraphs[0].font.size = Pt(11)
        c1.text_frame.paragraphs[0].font.color.rgb = primary_color
        
        # Narrative
        c2 = table.cell(r_idx, 2)
        c2.text_frame.paragraphs[0].text = desc_val
        c2.text_frame.paragraphs[0].font.size = Pt(11)
        c2.text_frame.paragraphs[0].font.color.rgb = text_light
        
        # Dark style for cells
        for c in range(3):
            cell = table.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    # Save presentation
    output_path = "c:\\Users\\HP\\OneDrive\\Desktop\\SOCF2.0\\PaperPilot_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
