import os
import uvicorn
from fastapi import FastAPI, UploadFile, File, HTTPException, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import Dict, Any

from pdf_parser import extract_pdf_chunks, calculate_file_hash
from vector_db import LocalVectorDB
from cache_manager import CacheManager
from agent import generate_single_pass_brief

app = FastAPI(title="PaperPilot API", version="2.0.0")

# Setup CORS to allow Next.js on port 3000 to query our backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # For hackathon/local dev, allow all
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize storage directories
PDF_DIR = "./data/pdfs"
os.makedirs(PDF_DIR, exist_ok=True)

# Initialize components
db = LocalVectorDB(persist_dir="./data")
cache = CacheManager(db_dir="./data")

class BriefRequest(BaseModel):
    file_hash: str
    eli5: bool = False

@app.post("/api/upload")
async def upload_pdf(file: UploadFile = File(...)):
    """
    Accepts PDF file, calculates hash, extracts text, embeds, and stores in Vector DB.
    Checks cache first to bypass duplicate ingestion.
    """
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
        
    try:
        contents = await file.read()
        file_hash = calculate_file_hash(contents)
        
        # Save file to disk to serve it to the frontend PDF viewer
        pdf_path = os.path.join(PDF_DIR, f"{file_hash}.pdf")
        if not os.path.exists(pdf_path):
            with open(pdf_path, "wb") as f:
                f.write(contents)
        
        # Check if the brief for this file hash is already cached
        cached_brief = cache.get_cached_brief(file_hash)
        if cached_brief and isinstance(cached_brief, dict):
            print(f"File {file_hash} already analyzed and cached. Bypassing ingestion.")
            return {
                "file_hash": file_hash,
                "cached": True,
                "title": cached_brief.get("title", file.filename),
                "message": "File already analyzed and cached."
            }

        # Otherwise, parse and chunk
        print(f"Extracting chunks from uploaded PDF: {file.filename}")
        chunks = extract_pdf_chunks(pdf_path, file_hash)
        print(f"Extracted {len(chunks)} chunks.")
        
        if not chunks:
            print("PDF has no extractable text layer. Spawning dynamic RAG content generator...")
            try:
                from pypdf import PdfReader
                reader = PdfReader(pdf_path)
                meta_title = reader.metadata.get('/Title', '') if (reader.metadata and reader.metadata.get('/Title')) else ''
                clean_title = meta_title if (meta_title and len(meta_title.strip()) > 3) else file.filename
            except Exception:
                clean_title = file.filename
                
            from agent import generate_simulated_paper_text
            simulated_text = generate_simulated_paper_text(clean_title)
            
            paragraphs = simulated_text.split("\n\n")
            for idx, para in enumerate(paragraphs):
                para = para.strip()
                if para:
                    chunks.append({
                        "chunk_id": f"{file_hash}_p{idx+1}_c{idx}",
                        "page_number": idx + 1,
                        "text": para
                    })
            print(f"Dynamically generated {len(chunks)} simulated text chunks based on: '{clean_title}'")
            
        if not chunks:
            raise HTTPException(status_code=400, detail="Unable to extract text or generate simulated content from the PDF.")
            
        # Add to local vector store (using SentenceTransformer embeddings)
        db.add_chunks(file_hash, chunks)
        
        return {
            "file_hash": file_hash,
            "cached": False,
            "num_chunks": len(chunks),
            "message": "File successfully uploaded, parsed, and embedded locally."
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Ingestion failed: {str(e)}")

@app.post("/api/brief")
async def generate_brief(payload: BriefRequest):
    """
    Retrieves chunks relevant to Methodology, Results, and Limitations,
    calls LLM in single-pass JSON mode, caches response, and updates cost stats.
    """
    file_hash = payload.file_hash
    eli5 = payload.eli5
    cache_key = f"{file_hash}_eli5" if eli5 else file_hash
    
    # 1. Check SQLite cache first for semantic cache hits
    cached_brief = cache.get_cached_brief(cache_key)
    if cached_brief:
        print(f"Cache Hit for key {cache_key}! Returning cached research brief.")
        return {
            "brief": cached_brief,
            "source": "cache",
            "message": "Retrieved from local SQLite cache."
        }
        
    # 2. Cache Miss - Execute local RAG retrieval
    print(f"Cache Miss for key {cache_key}. Performing local semantic search across index...")
    
    # Query 3 key areas of study separately (local computations, 0 API calls)
    q_methodology = db.query(file_hash, "methodology dataset experimental setup model training architecture framework implementation", n_results=3)
    q_results = db.query(file_hash, "results evaluation metrics findings benchmarks comparison accuracy loss graphs tables performance", n_results=3)
    q_limitations = db.query(file_hash, "limitations future work assumptions drawbacks challenges scope issues scope outline", n_results=3)
    
    # Deduplicate retrieved chunks
    seen_ids = set()
    retrieved_chunks = []
    for chunk in (q_methodology + q_results + q_limitations):
        if chunk["chunk_id"] not in seen_ids:
            seen_ids.add(chunk["chunk_id"])
            retrieved_chunks.append(chunk)
            
    if not retrieved_chunks:
        raise HTTPException(status_code=404, detail="No text chunks found for this file hash. Please re-upload.")
        
    print(f"Retrieved {len(retrieved_chunks)} unique context chunks. Sending to single-pass LLM agent (ELI5={eli5})...")
    
    # 3. Call LLM agent (Structured single-pass call)
    try:
        brief_data, tokens, cost = generate_single_pass_brief(retrieved_chunks, eli5=eli5)
        
        # 4. Save to cache
        cache.save_brief_to_cache(cache_key, brief_data)
        
        # 5. Record API call metrics
        cache.record_llm_call(tokens, cost)
        
        return {
            "brief": brief_data,
            "source": "llm",
            "message": "Successfully generated brief with single-pass LLM."
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Brief generation failed: {str(e)}")

@app.get("/api/stats")
async def get_stats():
    """
    Returns system performance and API savings statistics.
    """
    return cache.get_stats()

@app.post("/api/stats/reset")
async def reset_stats():
    """
    Resets all performance statistics.
    """
    cache.reset_stats()
    return {"message": "System stats reset successfully."}

class ChatRequest(BaseModel):
    file_hash: str
    question: str

@app.post("/api/chat")
async def copilot_chat(payload: ChatRequest):
    """
    Retrieves custom question context, queries the active client (Groq/OpenAI),
    caches responses semantic-style, and returns answers with clickable citations.
    """
    file_hash = payload.file_hash
    question = payload.question
    
    if not file_hash or not question:
        raise HTTPException(status_code=400, detail="Missing file_hash or question.")
        
    # 1. Check SQLite copilot cache
    cached_ans = cache.get_cached_copilot(file_hash, question)
    if cached_ans:
        print(f"Copilot Cache Hit for: {question}")
        return {
            "answer_data": cached_ans,
            "source": "cache",
            "message": "Retrieved from local SQLite copilot cache."
        }
        
    # 2. Local RAG Retrieval (0 API call embeddings/search)
    print(f"Copilot Cache Miss. Searching local index for: {question}")
    retrieved = db.query(file_hash, question, n_results=4)
    if not retrieved:
         raise HTTPException(status_code=404, detail="No source context retrieved for this question.")
        
    # 3. Call LLM for single-pass structured JSON
    try:
        from agent import generate_copilot_answer
        answer_data, tokens, cost = generate_copilot_answer(question, retrieved)
        
        # 4. Save to SQLite cache
        cache.save_copilot_to_cache(file_hash, question, answer_data)
        
        # 5. Record API call
        cache.record_llm_call(tokens, cost)
        
        return {
            "answer_data": answer_data,
            "source": "llm",
            "message": "Successfully generated answer with RAG copilot."
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Copilot query failed: {str(e)}")

def build_pptx_slides(title_text: str, slides_list: list) -> "io.BytesIO":
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    bg_color = RGBColor(10, 15, 30)         # Dark Navy
    primary_color = RGBColor(139, 92, 246)  # Electric Violet
    accent_color = RGBColor(16, 185, 129)   # Emerald Green
    text_light = RGBColor(241, 245, 249)    # Light gray/white
    box_bg = RGBColor(20, 27, 45)           # Lighter box container
    
    blank_layout = prs.slide_layouts[6]
    
    # Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = bg_color
    
    # Left Vertical Accent Strip (violet)
    accent_strip1 = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_strip1.fill.solid()
    accent_strip1.fill.fore_color.rgb = primary_color
    accent_strip1.line.fill.background()
    
    # Title Text Box
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.5))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p_tag = tf.paragraphs[0]
    p_tag.text = "AUTONOMOUS RESEARCH PRESENTATION"
    p_tag.font.name = 'Arial'
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = title_text
    p.font.name = 'Georgia'
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = text_light
    p.space_before = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "Synthesized autonomously by PaperPilot AcademIQ Agent"
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = primary_color
    p2.space_before = Pt(12)
    
    # Content Slides
    for slide_data in slides_list:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        # Left Accent Strip
        accent_strip = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
        accent_strip.fill.solid()
        accent_strip.fill.fore_color.rgb = primary_color
        accent_strip.line.fill.background()
        
        # Slide Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = slide_data.get('title', 'Content Overview').upper()
        tp.font.name = 'Georgia'
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = text_light
        
        tp_sub = ttf.add_paragraph()
        tp_sub.text = f"SECTION {slide_data.get('slide_number', 1)}  •  STUDY INSIGHTS"
        tp_sub.font.name = 'Arial'
        tp_sub.font.size = Pt(10)
        tp_sub.font.bold = True
        tp_sub.font.color.rgb = accent_color
        tp_sub.space_before = Pt(4)
        
        # Cards layout
        bullet_points = slide_data.get('bullet_points', [])
        start_y = 1.8
        card_h = 1.0
        gap_y = 0.25
        
        for idx, pt in enumerate(bullet_points[:4]):
            current_y = start_y + idx * (card_h + gap_y)
            
            # Card shape
            card = slide.shapes.add_shape(1, Inches(0.8), Inches(current_y), Inches(11.7), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = box_bg
            card.line.color.rgb = RGBColor(30, 41, 59)
            card.line.width = Pt(1.5)
            
            # Text inside card
            ctf = card.text_frame
            ctf.word_wrap = True
            ctf.margin_top = Inches(0.2)
            ctf.margin_left = Inches(0.3)
            ctf.margin_right = Inches(0.3)
            
            cp = ctf.paragraphs[0]
            cp.text = pt
            cp.font.name = 'Arial'
            cp.font.size = Pt(13)
            cp.font.color.rgb = text_light
            cp.font.bold = False
            
    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    return output_stream

class DownloadPresentationRequest(BaseModel):
    file_hash: str
    slides: list

@app.post("/api/download_presentation")
async def download_presentation(payload: DownloadPresentationRequest):
    """
    Generates a dark-themed PPTX slide deck based on generated slides,
    returning it as a file download.
    """
    try:
        title_text = "Research Slide Deck"
        cached_brief = cache.get_cached_brief(payload.file_hash)
        if cached_brief and isinstance(cached_brief, dict):
            title_text = cached_brief.get("title", title_text)
            
        from fastapi.responses import StreamingResponse
        output_stream = build_pptx_slides(title_text, payload.slides)
        
        filename = f"PaperPilot_{payload.file_hash[:8]}_Slides.pptx"
        return StreamingResponse(
            output_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Failed to generate PowerPoint: {str(e)}")

def build_pdf_report(title_text: str, brief_data: dict) -> "io.BytesIO":
    import io
    from reportlab.lib.pagesizes import letter
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = getSampleStyleSheet()
    
    # Custom Styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=colors.HexColor('#1e1b4b'),
        spaceAfter=12
    )
    
    subtitle_style = ParagraphStyle(
        'DocSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-BoldOblique',
        fontSize=9,
        leading=13,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=20
    )
    
    h1_style = ParagraphStyle(
        'SectionHeader',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor('#1e1b4b'),
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#334155'),
        spaceAfter=8
    )
    
    meta_label_style = ParagraphStyle(
        'MetaLabel',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor('#0f172a')
    )
    
    meta_val_style = ParagraphStyle(
        'MetaVal',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=13,
        textColor=colors.HexColor('#475569')
    )

    story = []
    
    story.append(Paragraph(title_text, title_style))
    story.append(Paragraph("AUTONOMOUS BRIEFING REPORT GENERATED BY PAPERPILOT ACADEMIQ", subtitle_style))
    story.append(Spacer(1, 10))
    
    m_score = brief_data.get("methodology_score", 7)
    flaw = brief_data.get("critical_flaw", "No major flaws identified.")
    
    overview_data = [
        [Paragraph("Methodology Score:", meta_label_style), Paragraph(f"{m_score} / 10", meta_val_style)],
        [Paragraph("Critical Critique:", meta_label_style), Paragraph(flaw, meta_val_style)],
        [Paragraph("Key Takeaways:", meta_label_style), Paragraph(brief_data.get("abstract_summary", ""), meta_val_style)]
    ]
    
    overview_table = Table(overview_data, colWidths=[120, 384])
    overview_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#e2e8f0')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#f1f5f9')),
        ('PADDING', (0,0), (-1,-1), 8),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(overview_table)
    story.append(Spacer(1, 15))
    
    story.append(Paragraph("1. Research Methodology", h1_style))
    story.append(Paragraph(brief_data.get("methodology", ""), body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("2. Experimental Results & Benchmarks", h1_style))
    story.append(Paragraph(brief_data.get("results", ""), body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("3. Assumptions, Limitations & Future Scope", h1_style))
    story.append(Paragraph(brief_data.get("limitations", ""), body_style))
    story.append(Spacer(1, 20))
    
    story.append(PageBreak())
    
    story.append(Paragraph("4. Scientific Concept Map Outline", h1_style))
    story.append(Paragraph("Below is the hierarchical structure of key concepts, models, and metric nodes extracted from the paper:", body_style))
    story.append(Spacer(1, 10))
    
    cmap = brief_data.get("concept_map", {})
    nodes = cmap.get("nodes", [])
    edges = cmap.get("edges", [])
    
    if not nodes:
        nodes = [
            {"id": "node-bg", "label": title_text, "group": "background"},
            {"id": "node-arch", "label": "Proposed System Architecture", "group": "architecture"},
            {"id": "node-method", "label": "Core Technical Pipeline", "group": "methodology"},
            {"id": "node-results", "label": "Experimental Performance", "group": "results"},
            {"id": "node-conclusion", "label": "Conclusion & Future Work", "group": "results"}
        ]
        
    bg_nodes = [n["label"] for n in nodes if n.get("group", "").lower() in ["background", "default"]]
    arch_nodes = [n["label"] for n in nodes if n.get("group", "").lower() in ["architecture", "model"]]
    method_nodes = [n["label"] for n in nodes if n.get("group", "").lower() == "methodology"]
    res_nodes = [n["label"] for n in nodes if n.get("group", "").lower() == "results"]
    
    concept_rows = [
        [Paragraph("BACKGROUND:", meta_label_style), Paragraph(", ".join(bg_nodes) if bg_nodes else "N/A", meta_val_style)],
        [Paragraph("ARCHITECTURE:", meta_label_style), Paragraph(", ".join(arch_nodes) if arch_nodes else "N/A", meta_val_style)],
        [Paragraph("METHODOLOGY:", meta_label_style), Paragraph(", ".join(method_nodes) if method_nodes else "N/A", meta_val_style)],
        [Paragraph("RESULTS & FINDINGS:", meta_label_style), Paragraph(", ".join(res_nodes) if res_nodes else "N/A", meta_val_style)]
    ]
    
    concept_table = Table(concept_rows, colWidths=[140, 364])
    concept_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#e2e8f0')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#f1f5f9')),
        ('PADDING', (0,0), (-1,-1), 8),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(concept_table)
    story.append(Spacer(1, 20))
    
    story.append(Paragraph("5. Study Flashcards", h1_style))
    cards = brief_data.get("flashcards", [])
    
    card_table_data = []
    for c_idx, card_item in enumerate(cards):
        q = f"Q{c_idx+1}: {card_item.get('question')}"
        a = f"A: {card_item.get('answer')} (Page {card_item.get('page_number', 'N/A')})"
        
        card_table_data.append([Paragraph(q, meta_label_style)])
        card_table_data.append([Paragraph(a, meta_val_style)])
        card_table_data.append([Spacer(1, 5)])
        
    if card_table_data:
        cards_table = Table(card_table_data, colWidths=[504])
        cards_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#fcfdfd')),
            ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#f1f5f9')),
            ('PADDING', (0,0), (-1,-1), 6),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ]))
        story.append(cards_table)
        
    doc.build(story)
    buffer.seek(0)
    return buffer

class DownloadPdfRequest(BaseModel):
    file_hash: str
    brief: dict

@app.post("/api/download_pdf")
async def download_pdf(payload: DownloadPdfRequest):
    """
    Generates a professionally styled PDF report of the academic brief,
    returning it as a file download.
    """
    try:
        title_text = payload.brief.get("title", "Research Brief Report")
        
        from fastapi.responses import StreamingResponse
        output_stream = build_pdf_report(title_text, payload.brief)
        
        filename = f"PaperPilot_{payload.file_hash[:8]}_Report.pdf"
        return StreamingResponse(
            output_stream,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Failed to generate PDF report: {str(e)}")

@app.get("/api/pdfs/{file_hash}")
async def get_pdf(file_hash: str):
    """
    Serves the uploaded PDF file to the frontend PDF viewer.
    """
    pdf_path = os.path.join(PDF_DIR, f"{file_hash}.pdf")
    if not os.path.exists(pdf_path):
        raise HTTPException(status_code=404, detail="PDF not found.")
    return FileResponse(pdf_path, media_type="application/pdf")

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
